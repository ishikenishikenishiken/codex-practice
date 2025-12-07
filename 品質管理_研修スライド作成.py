import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Cm
from PIL import Image

def create_dummy_image(path, text, size=(640, 480)):
    img = Image.new('RGB', size, color=(200, 200, 200))
    from PIL import ImageDraw
    draw = ImageDraw.Draw(img)
    textxy = (size[0] // 10, size[1] // 2)
    draw.text(textxy, text, fill=(50, 50, 50))
    img.save(path)

def main():
    prs = Presentation()

    # タイトルスライド
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "品質管理基礎研修"
    slide.placeholders[1].text = "2025年12月"

    # 目次
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "目次"
    content = [
        "1. 品質管理とは",
        "2. 品質管理の重要性",
        "3. QC七つ道具",
        "4. 品質改善事例",
        "5. 今後の課題"
    ]
    tf = slide.placeholders[1].text_frame
    for item in content:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
    tf.paragraphs[0].text = content[0]  # 最初だけ置換

    # 各ページ（内容5ページ）
    pages = [
        {"title": "品質管理とは", "body": "品質管理は、製品やサービスの品質を維持・向上させるための活動を指します。"},
        {"title": "品質管理の重要性", "body": "品質管理なくして顧客満足なし。企業の信頼構築、コスト削減、クレーム防止につながります。"},
        {"title": "QC七つ道具", "body": "特性要因図、パレート図、ヒストグラムなど、データを活用した問題解決ツール。"},
        {"title": "品質改善事例", "body": "過去の実際の事例から、改善活動の成果や方法を紹介。"},
        {"title": "今後の課題", "body": "AI活用やグローバル展開など、進化する品質管理への備えが求められます。"}
    ]
    img_dir = "dummy_imgs"
    os.makedirs(img_dir, exist_ok=True)
    for i, page in enumerate(pages):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = page["title"]
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = page['body']

        # 仮画像作成＆挿入
        img_path = f"{img_dir}/img_{i+1}.png"
        create_dummy_image(img_path, f"仮画像 {i+1}")
        pic_left = Cm(8)
        pic_top = Cm(5)
        pic_width = Cm(7)
        slide.shapes.add_picture(img_path, pic_left, pic_top, width=pic_width)

    prs.save("品質管理_研修スライド.pptx")
    print("品質管理_研修スライド.pptx を作成しました")

if __name__ == "__main__":
    main()

